####################################################################################
# Preprocessors for Text Extraction
# by JW
#
# A powerful collection of utilities to extract text from websites, PDFs, and other
# file formats.
# (see https://github.com/Gladiator07/Text-Summarizer/blob/main/utils.py)
# 
# preprocessors / text_extraction.py
####################################################################################

# IMPORT STATEMENTS ----------------------------------------------------------------
import re
import requests
import docx2txt
from io import StringIO
from PyPDF2 import PdfReader
from pptx import Presentation

from bs4 import BeautifulSoup
from nltk.tokenize import sent_tokenize

emoji_pattern = re.compile(
    "["
    u"\U0001F600-\U0001F64F"  # emoticons
    u"\U0001F300-\U0001F5FF"  # symbols & pictographs
    u"\U0001F680-\U0001F6FF"  # transport & map symbols
    u"\U0001F1E0-\U0001F1FF"  # flags (iOS)
    u"\U00002702-\U000027B0"
    u"\U000024C2-\U0001F251"
    "]+",
    flags=re.UNICODE,
)


def clean_text(x):
    # x = x.lower()  # lowercase
    x = x.encode("ascii", "ignore").decode()  # unicode
    x = re.sub(r"https*\S+", " ", x)  # url
    x = re.sub(r"@\S+", " ", x)  # mentions
    x = re.sub(r"#\S+", " ", x)  # hastags
    # x = x.replace("'", "")  # remove ticks
    # x = re.sub("[%s]" % re.escape(string.punctuation), " ", x)  # punctuation
    x = re.sub(r"\w*\d+\w*", "", x)  # numbers
    x = re.sub(r"\s{2,}", " ", x)  # over spaces
    x = emoji_pattern.sub(r"", x)  # emojis
    x = re.sub("[^.,!?A-Za-z0-9]+", " ", x)  # special charachters except .,!?
    # x = re.sub(r"(\w)([A-Z])", r"\1 \2", x) # Insert spaces before capital letters in order to split words

    # ' '.join(word for word in x.split() if len(word)>1)
    # ' '.join(word for word in x.split() if len(word)<20)

    return x

def format_url(url):
    if not re.match('(?:http|ftp|https)://', url):
        return 'http://{}'.format(url)
    return url

def fetch_article_text(url: str):

    r = requests.get(url)
    soup = BeautifulSoup(r.text, "html.parser")
    results = soup.find_all(["h1", "p"])
    text = [result.text for result in results]
    ARTICLE = " ".join(text)
    ARTICLE = ARTICLE.replace(".", ".<eos>")
    ARTICLE = ARTICLE.replace("!", "!<eos>")
    ARTICLE = ARTICLE.replace("?", "?<eos>")
    sentences = ARTICLE.split("<eos>")
    # current_chunk = 0
    # chunks = []
    # for sentence in sentences:
    #     if len(chunks) == current_chunk + 1:
    #         if len(chunks[current_chunk]) + len(sentence.split(" ")) <= 500:
    #             chunks[current_chunk].extend(sentence.split(" "))
    #         else:
    #             current_chunk += 1
    #             chunks.append(sentence.split(" "))
    #     else:
    #         chunks.append(sentence.split(" "))

    # for chunk_id in range(len(chunks)):
    #     chunks[chunk_id] = " ".join(chunks[chunk_id])

    for sentence in sentences:
        sentence = clean_text(sentence)

    return ARTICLE, sentences


def preprocess_text_for_abstractive_summarization(tokenizer, text):
    sentences = sent_tokenize(text)

    # initialize
    length = 0
    chunk = ""
    chunks = []
    count = -1
    for sentence in sentences:
        count += 1
        combined_length = (
            len(tokenizer.tokenize(sentence)) + length
        )  # add the no. of sentence tokens to the length counter

        if combined_length <= tokenizer.max_len_single_sentence:  # if it doesn't exceed
            chunk += sentence + " "  # add the sentence to the chunk
            length = combined_length  # update the length counter

            # if it is the last sentence
            if count == len(sentences) - 1:
                chunks.append(chunk.strip())  # save the chunk

        else:
            chunks.append(chunk.strip())  # save the chunk

            # reset
            length = 0
            chunk = ""

            # take care of the overflow sentence
            chunk += sentence + " "
            length = len(tokenizer.tokenize(sentence))

    return chunks


def read_pdf(file):
    pdfReader = PdfReader(file)
    count = len(pdfReader.pages)
    all_page_text = ""
    for i in range(count):
        page = pdfReader.pages[i]
        all_page_text += page.extract_text()

    return all_page_text


def read_text_from_file(file):

    print(file.type)

    # read text file
    if file.type == "text/plain" or file.type == "text/markdown":
        # To convert to a string based IO:
        stringio = StringIO(file.getvalue().decode("utf-8"))

        # To read file as string:
        file_content = stringio.read()

    # read pdf file
    elif file.type == "application/pdf":
        file_content = read_pdf(file)

    # read docx file
    elif (
        file.type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        file_content = docx2txt.process(file)

    elif (
        file.type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        file_content = ""
        pres = Presentation(file)
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    file_content = file_content + " " + shape.text

    else:
        file_content = ""

    return file_content